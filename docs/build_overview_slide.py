"""Generates docs/SQL-Project-Overview.pptx — a single overview slide."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

DARK = RGBColor(0x24, 0x24, 0x24)
GREY = RGBColor(0x60, 0x60, 0x60)
LIGHT = RGBColor(0xF3, 0xF3, 0xF3)
LINE = RGBColor(0xD9, 0xD9, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

pillars = [
    ("Schema as Code", RGBColor(0x0F, 0x6C, 0xBD), [
        "SDK-style (Microsoft.Build.Sql)",
        "One object per file, auto-globbed",
        "Organized by schema / object type",
        "Declarative — no hand-written ALTERs",
    ]),
    ("Author", RGBColor(0x10, 0x7C, 0x10), [
        "New empty project",
        "Import from existing database",
        "Schema Compare (drift detection)",
        "VS Code SQL Database Projects ext.",
    ]),
    ("Build & Validate", RGBColor(0x5C, 0x2D, 0x91), [
        "dotnet build → DACPAC",
        "Azure SQL surface validation (DSP)",
        "Warnings-as-errors (Release)",
        "Static code analysis (SR rules)",
    ]),
    ("Deploy", RGBColor(0xC2, 0x39, 0xB3), [
        "SqlPackage / DACPAC publish",
        "SqlCmd variables per environment",
        "Pre / Post-deployment scripts",
        "BlockOnPossibleDataLoss",
    ]),
    ("CI/CD", RGBColor(0xD8, 0x3B, 0x01), [
        "Azure DevOps + GitHub Actions",
        "OIDC / WIF — no stored secrets",
        "Native SqlAzureDacpacDeployment",
        "Automated smoke tests",
    ]),
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

SW = 13.333

# accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SW), Inches(0.18))
bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0x0F, 0x6C, 0xBD); bar.line.fill.background()

# title
tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.35), Inches(SW - 0.8), Inches(0.9))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "SQL Database Project — Essential Functionality & Tools"
p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = DARK

# subtitle
sb = slide.shapes.add_textbox(Inches(0.4), Inches(1.10), Inches(SW - 0.8), Inches(0.45))
sf = sb.text_frame; sf.word_wrap = True
sp = sf.paragraphs[0]
sp.text = "Database as code on Azure SQL — author, build, compare, and deploy one portable DACPAC"
sp.font.size = Pt(15); sp.font.color.rgb = GREY

# flow banner (.sqlproj → build → .dacpac → deploy → Database)
img_path = os.path.join(os.path.dirname(__file__), "build-deploy.png")
img_w = 8.4
img_h = img_w * 472 / 1800
slide.shapes.add_picture(img_path, Inches((SW - img_w) / 2), Inches(1.62), width=Inches(img_w))

# cards
left0 = 0.4
card_w = 2.34
gap = 0.20
top = 3.98
card_h = 2.72
head_h = 0.50

for i, (name, color, items) in enumerate(pillars):
    x = left0 + i * (card_w + gap)
    # card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(card_w), Inches(card_h))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LINE; card.line.width = Pt(1)
    card.shadow.inherit = False
    # header band
    head = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(card_w), Inches(head_h))
    head.fill.solid(); head.fill.fore_color.rgb = color; head.line.fill.background()
    head.shadow.inherit = False
    htf = head.text_frame; htf.word_wrap = True
    htf.margin_top = Pt(2); htf.margin_bottom = Pt(2)
    hp = htf.paragraphs[0]; hp.alignment = PP_ALIGN.CENTER
    hp.text = name; hp.font.size = Pt(15); hp.font.bold = True; hp.font.color.rgb = WHITE
    htf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # body bullets
    body = slide.shapes.add_textbox(Inches(x + 0.12), Inches(top + head_h + 0.12), Inches(card_w - 0.24), Inches(card_h - head_h - 0.24))
    btf = body.text_frame; btf.word_wrap = True
    for j, item in enumerate(items):
        para = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
        para.text = "•  " + item
        para.font.size = Pt(11); para.font.color.rgb = DARK
        para.space_after = Pt(6)

# footer
ft = slide.shapes.add_textbox(Inches(0.4), Inches(6.92), Inches(SW - 0.8), Inches(0.4))
ftf = ft.text_frame
fp = ftf.paragraphs[0]
fp.text = "One repo, one artifact → validated build → repeatable, safe deploy to every environment"
fp.font.size = Pt(12); fp.font.italic = True; fp.font.color.rgb = GREY

out = os.path.join(os.path.dirname(__file__), "SQL-Project-Overview.pptx")
prs.save(out)
print("saved", out)
